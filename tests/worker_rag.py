#!/usr/bin/env python3
import json, os, tempfile
os.environ["MODELRIG_DB"] = tempfile.mktemp(suffix=".db")

# Stub embeddings BEFORE any request: 26-dim lowercase letter-count vector so
# cosine is meaningful and the whole pipeline runs without a real Ollama.
import app.ollama_client as oc

def _vec(text: str):
    v = [0.0] * 26
    for ch in text.lower():
        i = ord(ch) - 97
        if 0 <= i < 26:
            v[i] += 1.0
    return v

async def fake_embed(text, model=None):
    return _vec(text)

oc.embed = fake_embed  # rag calls oc.embed at runtime → picks this up


async def fake_chat_stream(messages, model=None):
    for c in ["ctx", "-", "ans"]:
        yield (json.dumps({"message": {"content": c}}) + "\n").encode()

oc.chat_stream = fake_chat_stream

from fastapi.testclient import TestClient
from app.main import app
from app.rag import chunk_text

client = TestClient(app)
passed = failed = 0
def check(cond, name):
    global passed, failed
    print(f"  {'PASS' if cond else 'FAIL'}: {name}")
    passed += bool(cond); failed += (not cond)

# ---- chunk_text unit tests ----
check(chunk_text("") == [], "chunk empty -> []")
check(chunk_text("short text") == ["short text"], "chunk short -> single chunk")

long = " ".join(f"word{i}" for i in range(400))  # ~2.7k chars
cs = chunk_text(long, chunk_size=800, overlap=150)
check(len(cs) >= 3, f"chunk long -> multiple chunks ({len(cs)})")
check(all(len(c) <= 800 for c in cs), "every chunk within chunk_size")
check(set(long.split()) <= set(" ".join(cs).split()), "chunking loses no words")

# ---- ingest (chunked) via HTTP ----
r = client.post("/rag/ingest", json={"documents": [{"text": long, "source": "big"}],
                                     "chunk_size": 800, "overlap": 150})
check(r.status_code == 200, "ingest -> 200")
body = r.json()
check(body["chunks_added"] == len(cs), f"chunks_added matches chunk_text ({body['chunks_added']} vs {len(cs)})")
check(body["total"] == len(cs), "store total == chunks stored")

# ---- retrieval picks the nearest source ----
client.post("/rag/ingest", json={"documents": [
    {"text": "alpha bravo charlie delta", "source": "A"},
    {"text": "xray yankee zulu omega", "source": "B"}]})
r = client.post("/rag/query", json={"query": "yankee zulu", "top_k": 1, "synthesize": False})
check(r.status_code == 200, "query -> 200")
matches = r.json()["matches"]
check(len(matches) == 1 and matches[0]["source"] == "B",
      f"retrieval returns nearest source (got {matches[0]['source'] if matches else None})")
check("chunk_index" in matches[0] and "score" in matches[0], "match carries chunk_index + score")

# ---- source management: stats / sources / delete ----
# state now: "big" (len(cs) chunks), "A" (1), "B" (1)
r = client.get("/rag/stats")
st = r.json()
check(r.status_code == 200 and st["sources"] == 3, f"stats sources == 3 (got {st.get('sources')})")
check(st["chunks"] == len(cs) + 2, f"stats chunks == {len(cs)+2} (got {st.get('chunks')})")

r = client.get("/rag/sources")
srcs = {s["source"]: s["chunks"] for s in r.json()["sources"]}
check(r.status_code == 200 and set(srcs) == {"big", "A", "B"}, f"sources lists all (got {set(srcs)})")
check(srcs.get("big") == len(cs) and srcs.get("A") == 1, "per-source chunk counts correct")

r = client.delete("/rag/source", params={"source": "A"})
check(r.status_code == 200 and r.json()["removed"] == 1, "delete source A -> removed 1")
check(r.json()["total"] == len(cs) + 1, "total drops after delete")

r = client.get("/rag/stats")
check(r.json()["sources"] == 2, "stats sources == 2 after delete")

r = client.delete("/rag/source", params={"source": "does-not-exist"})
check(r.status_code == 404, "delete unknown source -> 404")

r = client.post("/rag/query", json={"query": "alpha bravo", "top_k": 3, "synthesize": False})
returned_sources = {m["source"] for m in r.json()["matches"]}
check("A" not in returned_sources, "deleted source no longer retrievable")

# ---- query restricted to a single source ----
# remaining state: "big" (many chunks), "B" (1)
r = client.post("/rag/query", json={"query": "xray yankee zulu", "top_k": 5,
                                     "synthesize": False, "source": "B"})
srcs = {m["source"] for m in r.json()["matches"]}
check(srcs == {"B"}, f"source filter returns only that source (got {srcs})")
r = client.post("/rag/query", json={"query": "word5 word6", "top_k": 5,
                                     "synthesize": False, "source": "big"})
srcs = {m["source"] for m in r.json()["matches"]}
check(srcs <= {"big"} and srcs, f"source filter 'big' -> only big (got {srcs})")

# ---- streaming RAG chat: first line = sources, rest = answer deltas ----
r = client.post("/rag/chat", json={"query": "xray yankee zulu", "top_k": 2})
check(r.status_code == 200, "rag-chat -> 200")
lines = [l for l in r.text.splitlines() if l.strip()]
head = json.loads(lines[0])
check("sources" in head and len(head["sources"]) >= 1, "rag-chat first line carries sources")
answer = "".join(json.loads(l).get("message", {}).get("content", "") for l in lines[1:])
check(answer == "ctx-ans", f"rag-chat streams the reassembled answer (got {answer!r})")

# ---- source chips are deduped per source name, not per chunk ----
# Ingest a doc long enough to split into multiple chunks, all sharing one
# source name, then confirm the rag/chat head emits ONE entry for that
# source (with a chunks count), not one-per-chunk. This is the exact
# on-device dup ("test.txt" shown twice for a 2-chunk file).
long_text = ("kanariefugl " * 60).strip()  # ~700 chars -> multiple chunks
client.post("/rag/ingest", json={"documents": [
    {"text": long_text, "source": "dedup-doc"}], "chunk_size": 200, "overlap": 40})
r = client.post("/rag/chat", json={"query": "kanariefugl", "top_k": 10, "source": "dedup-doc"})
head = json.loads(r.text.strip().split("\n")[0])
dedup_entries = [s for s in head["sources"] if s["source"] == "dedup-doc"]
check(len(dedup_entries) == 1,
      f"multiple chunks from one source -> exactly ONE source chip (got {len(dedup_entries)}: {head['sources']})")
check(dedup_entries and dedup_entries[0].get("chunks", 0) >= 2,
      f"deduped source entry carries a chunks count >= 2 (got {dedup_entries})")

# ---- min_score: irrelevant matches get filtered, not padded into top_k ----
# fake_embed is a pure 26-dim letter-count vector, so a query using only a
# letter absent from the stored text gives an exact, deterministic 0.0 cosine
# -- no ambiguity about whether the threshold is doing the filtering.
client.post("/rag/ingest", json={"documents": [
    {"text": "aeiou aeiou aeiou aeiou", "source": "vowels"}]})
r = client.post("/rag/query", json={"query": "zzzz", "top_k": 5, "synthesize": False,
                                     "min_score": 0.3, "source": "vowels"})
check(r.status_code == 200, "min_score query -> 200")
check(r.json()["matches"] == [],
      f"min_score=0.3 filters a zero-similarity match instead of padding top_k (got {r.json()['matches']})")

# same query, min_score=0.0 -- the same match now returns, proving the empty
# result above was the threshold filtering and not some unrelated bug
# (e.g. an empty store or a broken source filter)
r2 = client.post("/rag/query", json={"query": "zzzz", "top_k": 5, "synthesize": False,
                                      "min_score": 0.0, "source": "vowels"})
check(r2.status_code == 200 and len(r2.json()["matches"]) == 1,
      f"min_score=0.0 -> the same match returns, confirming it's the threshold (got {r2.json()['matches']})")

# empty-match synthesis: when min_score filters everything, /rag/query must
# still return an explicit answer (not omit the field and let the caller
# degrade to context-free chat -- the phone-vs-desktop divergence seen 7/7).
r3 = client.post("/rag/query", json={"query": "zzzz", "top_k": 5, "synthesize": True,
                                     "min_score": 0.3, "source": "vowels"})
check(r3.status_code == 200 and r3.json()["matches"] == [] and "answer" in r3.json()
      and r3.json()["answer"].strip() != "",
      f"empty matches still yields a non-empty answer field (got {r3.json()})")

# same for the streaming /rag/chat path: with no matches it must emit a chat
# delta (the don't-know message) and NOT call Ollama (which is dead here --
# if it tried, we'd see an error line instead of a clean message).
rc = client.post("/rag/chat", json={"query": "zzzz", "top_k": 5,
                                     "min_score": 0.3, "source": "vowels"})
lines_c = [json.loads(l) for l in rc.text.strip().splitlines()]
head_c = lines_c[0]
body_c = "".join(x.get("message", {}).get("content", "") for x in lines_c[1:])
has_error = any("error" in x for x in lines_c)
check(rc.status_code == 200 and head_c.get("sources") == [] and body_c.strip() != "" and not has_error,
      f"rag-chat with no matches streams a don't-know message without hitting Ollama (got {rc.text!r})")

# ---------------------------------------------------------------------------
# HTML ingest (stdlib -- always available, never 501)
# ---------------------------------------------------------------------------
from app import rag_html

check(rag_html.is_available(), "html: always available (stdlib)")

_h = (b"<html><head><title>Vadehavet &amp; tidevand</title>"
      b"<style>p{color:red}</style></head><body>"
      b"<nav>Forside Kontakt</nav><h1>Overskrift</h1><p>F&oslash;rste afsnit.</p>"
      b"<script>alert('x')</script>"
      b"<table><tr><td>VRAM</td><td>12 GB</td></tr></table>"
      b"<footer>Cookies</footer></body></html>")
_r = rag_html.extract_text(_h)
check(_r["title"] == "Vadehavet & tidevand", "html: title unescaped")
check("alert" not in _r["text"], "html: script contents dropped")
check("color:red" not in _r["text"], "html: style contents dropped")
check("Forside" not in _r["text"] and "Cookies" not in _r["text"],
      "html: nav/footer chrome dropped (would pollute every chunk)")
check("Første afsnit." in _r["text"], "html: danish entities decoded")
check("VRAM | 12 GB" in _r["text"], "html: table cells separated, not fused")
check("Overskrift" in _r["text"], "html: headings kept")

# Undecodable-as-utf8 bytes must still work (cp1252 fallback), not explode.
check("bl\u00e5" in rag_html.extract_text("<p>blå</p>".encode("cp1252"))["text"],
      "html: cp1252 fallback decodes danish letters")
try:
    rag_html.extract_text(b"<p></p>")
    _empty = rag_html.extract_text(b"<p></p>")["text"] == ""
except Exception:
    _empty = False
check(_empty, "html: empty document yields empty text (endpoint turns this into 422)")

# ---------------------------------------------------------------------------
# PPTX ingest (optional dependency, like PDF/DOCX)
# ---------------------------------------------------------------------------
from app import rag_pptx

if rag_pptx.is_available():
    import io as _io
    from pptx import Presentation as _P
    from pptx.util import Inches as _In
    _prs = _P()
    _s1 = _prs.slides.add_slide(_prs.slide_layouts[1])
    _s1.shapes.title.text = "Kaliv"
    _s1.placeholders[1].text = "Lokal AI"
    _s1.notes_slide.notes_text_frame.text = "Husk GPU'en"
    _s2 = _prs.slides.add_slide(_prs.slide_layouts[5])
    _t = _s2.shapes.add_table(1, 2, _In(1), _In(1), _In(4), _In(1)).table
    _t.cell(0, 0).text = "VRAM"; _t.cell(0, 1).text = "12 GB"
    _buf = _io.BytesIO(); _prs.save(_buf)
    _pr = rag_pptx.extract_text(_buf.getvalue())
    check(_pr["slides"] == 2, "pptx: slide count")
    check("[Slide 1]" in _pr["text"], "pptx: slide markers for traceability")
    check("Kaliv" in _pr["text"] and "Lokal AI" in _pr["text"], "pptx: title + body text")
    check("Husk GPU'en" in _pr["text"],
          "pptx: speaker notes captured (a deck of bullets embeds badly without them)")
    check("VRAM | 12 GB" in _pr["text"], "pptx: table cells")

    # Legacy binary .ppt must be rejected honestly, not with an opaque zip error.
    try:
        rag_pptx.extract_text(b"\xd0\xcf\x11\xe0" + b"0" * 32)
        _legacy = False
    except RuntimeError as e:
        _legacy = "legacy .ppt" in str(e)
    check(_legacy, "pptx: legacy .ppt rejected with a clear message")
else:
    print("  SKIP: python-pptx not installed (optional dependency)")

print(f"\n===== WORKER V1: {passed} passed, {failed} failed =====")
raise SystemExit(0 if failed == 0 else 1)
