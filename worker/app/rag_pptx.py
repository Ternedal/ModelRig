"""RAG PPTX ingest — text extraction with python-pptx.

Mirrors rag_docx.py: extracts text from an uploaded .pptx so it can go through
the EXISTING RAG ingest pipeline (chunk -> embed -> store). No new RAG logic.

What we pull, and why:
  - shape text (titles, bullets, text boxes) -- the obvious content
  - table cells -- real decks hide their numbers in tables
  - grouped shapes, recursively -- a grouped diagram is still content
  - speaker notes -- often the only prose in a slide deck; a deck of bullet
    fragments embeds badly, while the notes explain what the slide means

Slides are separated by a slide marker so the RAG chunker sees natural breaks
and a retrieved chunk can be traced back to a slide number.

OPTIONAL, like the PDF/DOCX/Voice backends: python-pptx is NOT a hard worker
dependency. Imported lazily; absent -> the endpoint returns a clean 501 with
install instructions, and the rest of the worker is unaffected.

Only the modern .pptx (OOXML) is supported. Legacy binary .ppt starts with the
OLE magic and is rejected with an honest message rather than an opaque zip
error. Fully testable without special hardware.
"""
from __future__ import annotations

import logging

import io


def is_available() -> bool:
    """True if python-pptx can be imported (installed)."""
    try:
        import pptx  # noqa: F401  (python-pptx)
        return True
    except Exception as exc:  # noqa: BLE001
        # Fail-closed is right: no python-pptx, no PPTX-indlæsning. Silence is not.
        #
        # This swallows more than ImportError, and on the rig it will: a broken
        # python-pptx wheel, a missing DLL, a CUDA mismatch. All of them arrive here
        # as "unavailable" with no reason, and the person reading the
        # capability list has to guess. F-501 was this exact shape -- an
        # `except Exception` hid an ImportError from a wrong class name for
        # eight releases, and the test passed because it asserted the failing
        # value and got it.
        logging.getLogger(__name__).info(
            "python-pptx er ikke tilgængelig (PPTX-indlæsning slået fra): %r", exc)
        return False


def _shape_texts(shapes) -> list[str]:
    """Text from a shape collection, descending into groups."""
    out: list[str] = []
    for shape in shapes:
        # A group is a container: its children hold the text, not the group.
        if getattr(shape, "shape_type", None) is not None and hasattr(shape, "shapes"):
            out.extend(_shape_texts(shape.shapes))
            continue
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = [(c.text or "").strip() for c in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    out.append(line)
            continue
        if getattr(shape, "has_text_frame", False):
            t = (shape.text_frame.text or "").strip()
            if t:
                out.append(t)
    return out


def extract_text(pptx_bytes: bytes) -> dict:
    """Extract text from .pptx bytes.

    Returns {text, slides, chars}. Raises RuntimeError (surfaced as 501/400 by
    the endpoint) on a missing backend or an unreadable/legacy file.
    """
    try:
        import pptx
    except Exception as e:
        raise RuntimeError(
            "PPTX ingest needs python-pptx. Install it on the rig with: pip install python-pptx"
        ) from e

    if pptx_bytes[:4] == b"\xd0\xcf\x11\xe0":
        raise RuntimeError("legacy .ppt is not supported; save as .pptx and retry")

    try:
        prs = pptx.Presentation(io.BytesIO(pptx_bytes))
    except Exception as e:
        raise RuntimeError(f"could not open PPTX: {e}") from e

    parts: list[str] = []
    slide_count = 0
    for i, slide in enumerate(prs.slides, start=1):
        slide_count += 1
        body = _shape_texts(slide.shapes)
        # Speaker notes carry the argument the bullets only gesture at.
        if slide.has_notes_slide:
            note = (slide.notes_slide.notes_text_frame.text or "").strip()
            if note:
                body.append(f"Noter: {note}")
        if body:
            parts.append(f"[Slide {i}]\n" + "\n".join(body))

    text = "\n\n".join(parts).strip()
    return {"text": text, "slides": slide_count, "chars": len(text)}
